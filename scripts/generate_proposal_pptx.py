"""
JOBTV サービス開発提案書 PowerPoint 自動生成スクリプト

Usage:
    pip install python-pptx
    python scripts/generate_proposal_pptx.py

Output:
    docs/presentations/jobtv_service_proposal.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Design tokens (shared with generate_executive_pptx.py)
# ---------------------------------------------------------------------------
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MEDIUM_GRAY = RGBColor(0x7F, 0x8C, 0x9A)
TEXT_DARK = RGBColor(0x2D, 0x3A, 0x4A)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
WARM_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
SUCCESS_GREEN = RGBColor(0x27, 0xAE, 0x60)
CARD_DARK_BG = RGBColor(0x24, 0x3B, 0x5E)

FONT_JP = "メイリオ"
FONT_EN = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "docs" / "presentations"
OUTPUT_FILE = OUTPUT_DIR / "jobtv_service_proposal.pptx"

TOTAL_SLIDES = 13


# ---------------------------------------------------------------------------
# Helpers (same as generate_executive_pptx.py)
# ---------------------------------------------------------------------------

def _set_font(run, size_pt, bold=False, color=TEXT_DARK, font_name=FONT_JP):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _add_bg_rect(slide, color=DARK_NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_accent_bar(slide, top=Inches(0), height=Inches(0.06)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_WIDTH, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def _add_footer(slide, page_num):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.95), Inches(12.13), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_GRAY
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {TOTAL_SLIDES}"
    _set_font(run, 9, color=MEDIUM_GRAY, font_name=FONT_EN)

    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(4), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "JOBTV Service Proposal  |  Confidential"
    _set_font(run2, 9, color=MEDIUM_GRAY, font_name=FONT_EN)


def _content_slide(prs, title_text, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.18), Inches(11), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_font(run, 26, bold=True, color=WHITE)

    _add_footer(slide, page_num)
    return slide


def _add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, font=FONT_JP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold=bold, color=color, font_name=font)
    return tb


def _add_bullet_list(slide, left, top, width, height, items, size=14, color=TEXT_DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = f"● {item}"
        _set_font(run, size, color=color)
    return tb


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.38 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            if r_idx == 0:
                _set_font(run, 11, bold=True, color=WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            else:
                _set_font(run, 11, color=TEXT_DARK)
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def _add_flow_boxes(slide, top, boxes, box_w=Inches(2.2), box_h=Inches(1.0)):
    n = len(boxes)
    gap = Inches(0.45)
    total_w = n * box_w + (n - 1) * gap
    start_x = int((SLIDE_WIDTH - total_w) / 2)

    for i, (label, desc) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE if i % 2 == 0 else DARK_NAVY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_font(run, 13, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        _set_font(run2, 10, color=RGBColor(0xCC, 0xDD, 0xEE))

        if i < n - 1:
            ax = x + box_w + Inches(0.03)
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, top + Inches(0.3), Inches(0.38), Inches(0.38)
            )
            a.fill.solid()
            a.fill.fore_color.rgb = ACCENT_BLUE
            a.line.fill.background()


def _add_card(slide, left, top, width, height, title, body, bg=LIGHT_BLUE_BG,
              title_color=DARK_NAVY, body_color=TEXT_DARK, body_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, 14, bold=True, color=title_color)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = body
    _set_font(run2, body_size, color=body_color)


def _placeholder_slide(prs, title_text, page_num, notice="後日差し替え予定"):
    """Create a placeholder slide with a centered notice."""
    slide = _content_slide(prs, title_text, page_num)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(2.5), Inches(6.3), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = notice
    _set_font(run, 28, bold=True, color=ACCENT_BLUE)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "本セクションの内容は別途ご相談のうえ記載いたします"
    _set_font(run2, 14, color=MEDIUM_GRAY)

    return slide


# ---------------------------------------------------------------------------
# Slide builders (13 slides)
# ---------------------------------------------------------------------------

def slide_01_cover(prs):
    """表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, DARK_NAVY)
    _add_accent_bar(slide, Inches(2.8), Inches(0.06))

    _add_textbox(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(1.2),
                 "JOBTV", 54, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)
    _add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(1.0),
                 "サービス開発のご提案", 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
                 "動画 × 新卒採用プラットフォーム", 20, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                 "2026年3月", 18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER, font=FONT_EN)
    _add_footer(slide, 1)


def slide_02_toc(prs):
    """目次"""
    slide = _content_slide(prs, "目次", 2)

    toc_items = [
        ("1.", "プロジェクト概要"),
        ("2.", "プロジェクトのスコープ"),
        ("3.", "主要機能一覧（JOBTV）"),
        ("4.", "主要機能一覧（Event System / Agent Manager）"),
        ("5.", "技術スタック"),
        ("6.", "システムアーキテクチャ"),
        ("7.", "セキュリティ・品質管理"),
        ("8.", "前提条件"),
        ("9.", "プロジェクト体制"),
        ("10.", "スケジュール"),
        ("11.", "お見積り"),
    ]

    for i, (num, title) in enumerate(toc_items):
        y = Inches(1.3) + Inches(0.48) * i
        # Number
        _add_textbox(slide, Inches(2.5), y, Inches(0.8), Inches(0.4),
                     num, 16, bold=True, color=ACCENT_BLUE, font=FONT_EN)
        # Title
        _add_textbox(slide, Inches(3.3), y, Inches(6), Inches(0.4),
                     title, 16, color=TEXT_DARK)
        # Dotted line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3.3), y + Inches(0.4), Inches(6), Pt(0.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_GRAY
        line.line.fill.background()


def slide_03_overview(prs):
    """1. プロジェクト概要"""
    slide = _content_slide(prs, "1. プロジェクト概要", 3)

    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
                 "動画で届ける、新しい新卒採用プラットフォーム", 22, bold=True, color=DARK_NAVY)

    _add_textbox(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.0),
                 "JOBTVは、企業が自社の魅力を動画で発信し、学生が「見て」「感じて」企業を選ぶ\n"
                 "新卒採用プラットフォームです。従来のテキスト中心の求人情報では伝わりにくい\n"
                 "職場の雰囲気や社員の生の声を、高品質な動画コンテンツで届けます。",
                 14, color=TEXT_DARK)

    # Concept cards
    concepts = [
        ("企業向け", "動画で企業文化・職場環境を\n発信。求人情報と連動した\nリッチなコンテンツ管理"),
        ("学生向け", "動画視聴で企業理解を深め、\n説明会予約・エントリーまで\nシームレスに完結"),
        ("運営向け", "コンテンツ審査・配信管理、\nイベント運営、候補者管理を\n一元化した管理画面"),
    ]
    for i, (title, body) in enumerate(concepts):
        x = Inches(0.8 + i * 4.1)
        _add_card(slide, x, Inches(3.3), Inches(3.7), Inches(2.2), title, body)

    # Vision
    _add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5),
                 "ビジョン:  企業と学生の「最適な出会い」を動画の力で創出する",
                 16, bold=True, color=ACCENT_BLUE)


def slide_04_scope(prs):
    """2. プロジェクトのスコープ"""
    slide = _content_slide(prs, "2. プロジェクトのスコープ", 4)

    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
                 "3つのアプリケーションで構成されるモノレポ構成", 18, bold=True, color=DARK_NAVY)

    apps = [
        ("JOBTV（メインサイト）",
         "学生向けサイト + 企業向け Studio\n"
         "● 動画視聴・企業ページ・求人閲覧\n"
         "● 企業によるコンテンツ入稿・管理\n"
         "● 管理者による審査・公開制御\n"
         "● LINE連携・通知・CMS",
         ACCENT_BLUE),
        ("Event System",
         "イベント運営プラットフォーム\n"
         "● 合同説明会の企画・運営\n"
         "● 参加企業・学生の管理\n"
         "● 予約・座席管理・当日運営\n"
         "● マッチングアルゴリズム",
         DARK_NAVY),
        ("Agent Manager",
         "候補者管理システム\n"
         "● 候補者情報の一元管理\n"
         "● 選考ステータス追跡\n"
         "● エージェント業務支援\n"
         "● レポート・分析機能",
         SUCCESS_GREEN),
    ]

    for i, (title, body, color) in enumerate(apps):
        x = Inches(0.8 + i * 4.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.1), Inches(3.7), Inches(3.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_font(run, 16, bold=True, color=color)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        run2 = p2.add_run()
        run2.text = body
        _set_font(run2, 12, color=TEXT_DARK)

    _add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
                 "※ 共通パッケージ（UI / DB型 / 設定）はモノレポ内で共有",
                 12, color=MEDIUM_GRAY)


def slide_05_features_jobtv(prs):
    """3. 主要機能一覧（JOBTV）"""
    slide = _content_slide(prs, "3. 主要機能一覧（JOBTV）", 5)

    rows = [
        ["カテゴリ", "機能", "概要"],
        ["動画配信", "HLSストリーミング", "AWS MediaConvert + CloudFront による高品質配信"],
        ["動画配信", "動画エディタ", "ブラウザ上でのトリミング・字幕編集"],
        ["企業ページ", "企業プロフィール", "動画・画像・テキストによるリッチな企業紹介"],
        ["求人", "求人管理", "Draft/Production ワークフローによる品質管理"],
        ["説明会", "オンライン説明会", "予約・リマインダー・Zoom連携"],
        ["LINE連携", "配信・CTA", "セグメント配信・リッチメニュー・スケジュール配信"],
        ["CMS", "LP・コンテンツ管理", "トップページ・特集ページの動的管理"],
        ["通知", "マルチチャネル", "Email / LINE / Slack / Sheets / アプリ内通知"],
        ["認証", "ロール管理", "Admin(MFA) / Recruiter / Candidate の3ロール"],
    ]

    _add_table(slide, Inches(0.8), Inches(1.3), Inches(11.7), rows,
               col_widths=[Inches(1.8), Inches(2.8), Inches(7.1)])


def slide_06_features_others(prs):
    """4. 主要機能一覧（Event System / Agent Manager）"""
    slide = _content_slide(prs, "4. 主要機能一覧（Event System / Agent Manager）", 6)

    # Event System
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Event System", 18, bold=True, color=ACCENT_BLUE)

    event_items = [
        "合同説明会イベントの企画・作成・公開",
        "参加企業の募集・管理・ブース割当",
        "学生の予約受付・座席管理",
        "当日の受付・出席管理",
        "イベント後のマッチング・フィードバック",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5),
                     event_items, size=13)

    # Agent Manager
    _add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Agent Manager", 18, bold=True, color=SUCCESS_GREEN)

    agent_items = [
        "候補者データベースの一元管理",
        "選考プロセスのステータス管理",
        "企業・候補者間コミュニケーション記録",
        "エージェント向けダッシュボード",
        "実績レポート・KPI分析",
    ]
    _add_bullet_list(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.5),
                     agent_items, size=13)

    # Separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.2), Pt(2), Inches(3.2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = LIGHT_GRAY
    sep.line.fill.background()

    # Shared capabilities
    _add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
                 "共通基盤", 16, bold=True, color=DARK_NAVY)

    shared = [
        ("認証・認可", "Supabase Auth + RLS\nロールベースアクセス制御"),
        ("通知", "Email / LINE / Slack\nマルチチャネル通知"),
        ("UI", "共通UIコンポーネント\nデザインシステム"),
        ("データ", "共通DB型定義\n型安全なクエリ"),
    ]
    for i, (title, body) in enumerate(shared):
        x = Inches(0.8 + i * 3.05)
        _add_card(slide, x, Inches(5.3), Inches(2.7), Inches(1.3), title, body, body_size=10)


def slide_07_tech_stack(prs):
    """5. 技術スタック"""
    slide = _content_slide(prs, "5. 技術スタック", 7)

    rows = [
        ["レイヤー", "技術", "用途"],
        ["フロントエンド", "Next.js 15 (App Router)", "SSR/SSG/RSC によるハイブリッドレンダリング"],
        ["言語", "TypeScript", "型安全な開発、フルスタック共通"],
        ["スタイリング", "Tailwind CSS + shadcn/ui", "ユーティリティファースト + コンポーネントライブラリ"],
        ["バックエンド", "Supabase (PostgreSQL)", "認証・DB・RLS・Realtime・Storage"],
        ["動画配信", "AWS (S3 + MediaConvert + CloudFront)", "HLS変換・CDN配信・署名付きURL"],
        ["ホスティング", "Vercel", "自動デプロイ・Edge Functions・Cron Jobs"],
        ["LINE連携", "LINE Messaging API", "リッチメニュー・Flex Message・配信"],
        ["メール", "Amazon SES / Resend", "トランザクションメール・通知"],
        ["モノレポ", "pnpm Workspaces + Turborepo", "依存管理・ビルド最適化"],
        ["E2Eテスト", "Playwright", "マルチアプリ対応・認証フィクスチャ"],
    ]

    _add_table(slide, Inches(0.8), Inches(1.3), Inches(11.7), rows,
               col_widths=[Inches(2.0), Inches(4.0), Inches(5.7)])


def slide_08_architecture(prs):
    """6. システムアーキテクチャ"""
    slide = _content_slide(prs, "6. システムアーキテクチャ", 8)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                 "モノレポ構成 × サーバーレスアーキテクチャ", 18, bold=True, color=DARK_NAVY)

    # Architecture flow
    _add_flow_boxes(slide, Inches(1.8), [
        ("Client", "ブラウザ / LINE"),
        ("Vercel", "Next.js SSR\nEdge Functions"),
        ("Supabase", "Auth / DB\nRLS / Storage"),
        ("AWS", "S3 / MediaConvert\nCloudFront"),
    ], box_w=Inches(2.5), box_h=Inches(1.2))

    # Monorepo structure
    _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.4),
                 "モノレポ構造", 16, bold=True, color=DARK_NAVY)

    mono_items = [
        "apps/jobtv — メインサイト (port 3000)",
        "apps/event-system — イベント管理 (port 3001)",
        "apps/agent-manager — 候補者管理 (port 3002)",
        "packages/ui — 共通UIコンポーネント",
        "packages/database — DB型定義・クライアント",
        "packages/config — 共通設定",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5),
                     mono_items, size=12)

    # Video pipeline
    _add_textbox(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(0.4),
                 "動画配信パイプライン", 16, bold=True, color=DARK_NAVY)

    video_items = [
        "1. 企業が動画をアップロード (S3)",
        "2. MediaConvert で HLS 変換",
        "3. CloudFront で CDN 配信",
        "4. 署名付き URL でセキュア再生",
        "5. アダプティブビットレート対応",
    ]
    _add_bullet_list(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(2.5),
                     video_items, size=12)


def slide_09_security(prs):
    """7. セキュリティ・品質管理"""
    slide = _content_slide(prs, "7. セキュリティ・品質管理", 9)

    # Security cards
    sec_items = [
        ("RLS（行レベルセキュリティ）",
         "全テーブルにRLSポリシーを適用。\nロール・テナントに基づくデータ分離を\nDB層で強制し、IDOR攻撃を防止。"),
        ("MFA（多要素認証）",
         "管理者アカウントにTOTPベースの\n多要素認証を必須化。\n不正アクセスリスクを低減。"),
        ("Draft/Production審査",
         "企業コンテンツはDraft→Production\nの2段階ワークフローで品質管理。\n管理者承認なしに公開されない。"),
        ("CAPTCHA・ボット対策",
         "Cloudflare Turnstileを導入。\nフォーム送信時のボット判定で\nスパム・不正登録を防止。"),
    ]
    for i, (title, body) in enumerate(sec_items):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.3) + (i // 2) * Inches(2.2)
        _add_card(slide, x, y, Inches(5.7), Inches(1.8), title, body, body_size=12)

    # Additional measures
    _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
                 "その他: XSS対策 / CSRF対策 / 署名付きURL / IP制限(admin) / 監査ログ / エラーバウンダリ",
                 13, color=MEDIUM_GRAY)


def slide_10_prerequisites(prs):
    """8. 前提条件"""
    slide = _content_slide(prs, "8. 前提条件", 10)

    # Contract
    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "契約形態・開発範囲", 18, bold=True, color=DARK_NAVY)

    contract_items = [
        "準委任契約を想定",
        "開発範囲は本提案書記載の機能を対象",
        "追加機能は別途お見積り",
        "開発環境・ステージング・本番の3環境構成",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.0),
                     contract_items, size=13)

    # Provided by client
    _add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "お客様にご準備いただくもの", 18, bold=True, color=DARK_NAVY)

    client_items = [
        "AWSアカウント（動画配信用）",
        "LINE公式アカウント（Messaging API）",
        "ドメイン・DNS管理",
        "動画素材・企業コンテンツ",
    ]
    _add_bullet_list(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.0),
                     client_items, size=13)

    # Exclusions
    _add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4),
                 "対象外事項", 18, bold=True, color=WARM_ORANGE)

    excl_items = [
        "動画撮影・編集（素材制作）",
        "ハードウェア・ネットワーク機器の調達",
        "既存システムからのデータ移行（別途相談）",
        "運用後の24/7監視体制（別途保守契約）",
    ]
    _add_bullet_list(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(1.8),
                     excl_items, size=13, color=TEXT_DARK)


def slide_11_team(prs):
    """9. プロジェクト体制（プレースホルダー）"""
    _placeholder_slide(prs, "9. プロジェクト体制", 11)


def slide_12_schedule(prs):
    """10. スケジュール（プレースホルダー）"""
    _placeholder_slide(prs, "10. スケジュール", 12)


def slide_13_estimate(prs):
    """11. お見積り（プレースホルダー）"""
    _placeholder_slide(prs, "11. お見積り", 13)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_overview(prs)
    slide_04_scope(prs)
    slide_05_features_jobtv(prs)
    slide_06_features_others(prs)
    slide_07_tech_stack(prs)
    slide_08_architecture(prs)
    slide_09_security(prs)
    slide_10_prerequisites(prs)
    slide_11_team(prs)
    slide_12_schedule(prs)
    slide_13_estimate(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"Generated: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
