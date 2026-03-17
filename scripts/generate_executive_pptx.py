"""
JOBTV エグゼクティブ向けプロジェクト報告 PowerPoint 自動生成スクリプト

Usage:
    pip install python-pptx
    python scripts/generate_executive_pptx.py

Output:
    docs/presentations/jobtv_executive_report.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MEDIUM_GRAY = RGBColor(0x7F, 0x8C, 0x9A)
TEXT_DARK = RGBColor(0x2D, 0x3A, 0x4A)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
WARM_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
SUCCESS_GREEN = RGBColor(0x27, 0xAE, 0x60)
CARD_DARK_BG = RGBColor(0x24, 0x3B, 0x5E)

FONT_JP = "メイリオ"
FONT_EN = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "docs" / "presentations"
OUTPUT_FILE = OUTPUT_DIR / "jobtv_executive_report.pptx"

TOTAL_SLIDES = 21


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_font(run, size_pt, bold=False, color=TEXT_DARK, font_name=FONT_JP):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _add_bg_rect(slide, color=DARK_NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_accent_bar(slide, top=Inches(0), height=Inches(0.06)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_WIDTH, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def _add_footer(slide, page_num):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.95), Inches(12.13), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_GRAY
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {TOTAL_SLIDES}"
    _set_font(run, 9, color=MEDIUM_GRAY, font_name=FONT_EN)

    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(4), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "JOBTV Project Report  |  Confidential"
    _set_font(run2, 9, color=MEDIUM_GRAY, font_name=FONT_EN)


def _content_slide(prs, title_text, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.18), Inches(11), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_font(run, 26, bold=True, color=WHITE)

    _add_footer(slide, page_num)
    return slide


def _add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, font=FONT_JP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold=bold, color=color, font_name=font)
    return tb


def _add_bullet_list(slide, left, top, width, height, items, size=14, color=TEXT_DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = f"● {item}"
        _set_font(run, size, color=color)
    return tb


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.38 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            if r_idx == 0:
                _set_font(run, 11, bold=True, color=WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            else:
                _set_font(run, 11, color=TEXT_DARK)
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def _add_flow_boxes(slide, top, boxes, box_w=Inches(2.2), box_h=Inches(1.0)):
    n = len(boxes)
    gap = Inches(0.45)
    total_w = n * box_w + (n - 1) * gap
    start_x = int((SLIDE_WIDTH - total_w) / 2)

    for i, (label, desc) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE if i % 2 == 0 else DARK_NAVY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_font(run, 13, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        _set_font(run2, 10, color=RGBColor(0xCC, 0xDD, 0xEE))

        if i < n - 1:
            ax = x + box_w + Inches(0.03)
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, ax, top + Inches(0.3), Inches(0.38), Inches(0.38)
            )
            a.fill.solid()
            a.fill.fore_color.rgb = ACCENT_BLUE
            a.line.fill.background()


def _add_card(slide, left, top, width, height, title, body, bg=LIGHT_BLUE_BG,
              title_color=DARK_NAVY, body_color=TEXT_DARK, body_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, 14, bold=True, color=title_color)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = body
    _set_font(run2, body_size, color=body_color)


def _add_icon_card(slide, left, top, width, height, icon_text, title, body,
                   icon_color=ACCENT_BLUE):
    """Card with a large icon/number on the left."""
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE_BG
    shape.line.color.rgb = RGBColor(0xD0, 0xE0, 0xF0)
    shape.line.width = Pt(1)

    # Icon circle
    circle_size = Inches(0.55)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cr = ctf.paragraphs[0].add_run()
    cr.text = icon_text
    _set_font(cr, 14, bold=True, color=WHITE, font_name=FONT_EN)

    # Title
    _add_textbox(slide, left + Inches(0.85), top + Inches(0.1), width - Inches(1.0), Inches(0.35),
                 title, 13, bold=True, color=DARK_NAVY)
    # Body
    _add_textbox(slide, left + Inches(0.85), top + Inches(0.42), width - Inches(1.0),
                 height - Inches(0.5), body, 11, color=MEDIUM_GRAY)


def _add_shield_row(slide, left, top, width, items):
    """Horizontal row of shield-style security items."""
    n = len(items)
    item_w = (width - Inches(0.2) * (n - 1)) / n
    for i, (label, detail) in enumerate(items):
        x = left + i * (item_w + Inches(0.2))
        # Shield shape
        shield = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, item_w, Inches(1.15)
        )
        shield.fill.solid()
        shield.fill.fore_color.rgb = LIGHT_BLUE_BG
        shield.line.color.rgb = ACCENT_BLUE
        shield.line.width = Pt(1.5)

        tf = shield.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _set_font(run, 12, bold=True, color=DARK_NAVY)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(3)
        run2 = p2.add_run()
        run2.text = detail
        _set_font(run2, 10, color=MEDIUM_GRAY)


# ---------------------------------------------------------------------------
# Slide builders (20 slides)
# ---------------------------------------------------------------------------


def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, DARK_NAVY)

    _add_accent_bar(slide, Inches(2.8), Inches(0.06))

    _add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
                 "JOBTV", 54, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)
    _add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.0),
                 "プロジェクト報告", 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
                 "動画 × 新卒採用プラットフォーム", 20, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                 "2026年3月", 18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER, font=FONT_EN)
    _add_footer(slide, 1)


def slide_02_exec_summary(prs):
    slide = _content_slide(prs, "エグゼクティブサマリー", 2)

    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
                 "企業文化を「見て」選ぶ時代へ", 22, bold=True, color=DARK_NAVY)

    # 3 large cards
    cards = [
        (ACCENT_BLUE, "動画で届ける",
         "HLS アダプティブストリーミングで\nあらゆる環境で快適視聴"),
        (DARK_NAVY, "品質と自走を両立",
         "Draft/Production 審査ワークフロー\n企業が入稿、管理者が承認・公開"),
        (SUCCESS_GREEN, "多チャネルで接点最大化",
         "Email / LINE / Slack / Sheets /\nアプリ内通知の 5 チャネル\nLINE スケジュール配信対応"),
    ]
    for i, (bg, title, body) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.1), Inches(3.8), Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_font(run, 18, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        run2 = p2.add_run()
        run2.text = body
        _set_font(run2, 13, color=RGBColor(0xDD, 0xEE, 0xFF))

    # Key tech badges
    badges = ["Next.js 16", "Supabase", "AWS (S3 / MediaConvert / CloudFront)", "Vercel", "TypeScript", "SEO / AIO"]
    badge_text = "    |    ".join(badges)
    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5),
                 badge_text, 13, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Stats row
    stats = [
        ("3", "ユーザーロール\nadmin / recruiter / candidate"),
        ("5", "通知チャネル\nEmail / LINE / Slack / Sheets / App"),
        ("5+", "審査対象コンテンツ\n企業ページ / 情報 / 求人 / 説明会 / 動画"),
        ("4", "定期処理（Cron）\nリマインダー / LINE 配信 / リトライ / 清掃"),
    ]
    for i, (num, desc) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)
        _add_textbox(slide, x, Inches(5.0), Inches(1.0), Inches(0.8),
                     num, 36, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)
        _add_textbox(slide, x + Inches(0.9), Inches(5.0), Inches(2.0), Inches(0.9),
                     desc, 11, color=MEDIUM_GRAY)


def slide_03_toc(prs):
    slide = _content_slide(prs, "目次", 3)

    sections = [
        ("1", "サービス概要", "全体像・ビジネスフロー"),
        ("2", "動画・コンテンツ管理", "HLS 配信・企業ページ・CMS"),
        ("3", "企業向け機能", "スタジオ・管理画面・審査フロー"),
        ("4", "技術基盤", "スタック・動画インフラ・DB 設計"),
        ("5", "セキュリティ", "多層防御・RLS・アクセス制御"),
        ("6", "運用・展望", "デプロイ・監視・ロードマップ"),
    ]

    top = Inches(1.6)
    for num, title, desc in sections:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), top, Inches(0.7), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        _set_font(run, 18, bold=True, color=WHITE, font_name=FONT_EN)

        _add_textbox(slide, Inches(3.0), top + Inches(0.05), Inches(3.5), Inches(0.4),
                     title, 18, bold=True, color=DARK_NAVY)
        _add_textbox(slide, Inches(6.8), top + Inches(0.1), Inches(5), Inches(0.4),
                     desc, 14, color=MEDIUM_GRAY)

        top += Inches(0.9)


def slide_04_overview(prs):
    slide = _content_slide(prs, "JOBTV とは", 4)

    _add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.5),
                 "新卒採用を動画で変える。企業密着・社員インタビュー・職場見学をリアルに届けるプラットフォーム。",
                 15, color=TEXT_DARK)

    # 3 role cards - visual
    roles = [
        ("学生", "Candidate", ACCENT_BLUE,
         "動画視聴・求人エントリー\n説明会予約・LINE 連携"),
        ("企業担当者", "Recruiter", DARK_NAVY,
         "動画アップロード・企業ページ\n求人・説明会の入稿管理"),
        ("管理者", "Admin", RGBColor(0x6C, 0x5C, 0xE7),
         "審査・承認・トップ CMS\nユーザー管理・LINE 配信"),
    ]
    for i, (label, role_en, bg, desc) in enumerate(roles):
        x = Inches(0.8 + i * 4.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.8), Inches(2.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        _set_font(run, 22, bold=True, color=WHITE)

        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = role_en
        _set_font(run2, 12, color=RGBColor(0xAA, 0xCC, 0xEE), font_name=FONT_EN)

        p3 = tf.add_paragraph()
        p3.space_before = Pt(12)
        run3 = p3.add_run()
        run3.text = desc
        _set_font(run3, 14, color=RGBColor(0xDD, 0xEE, 0xFF))

    _add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.3),
                 "本番ドメイン: https://media.jobtv.jp", 12, color=MEDIUM_GRAY, font=FONT_EN)


def slide_05_business_flow(prs):
    slide = _content_slide(prs, "ビジネスフロー", 5)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                 "企業側: 掲載から公開まで", 17, bold=True, color=DARK_NAVY)

    _add_flow_boxes(slide, Inches(1.7), [
        ("企業登録", "admin が\nアカウント作成"),
        ("コンテンツ入稿", "動画・企業ページ\n求人・説明会"),
        ("審査・承認", "admin が\nレビュー"),
        ("公開", "status = active\nで表示"),
    ])

    _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.4),
                 "学生側: 視聴からエントリーまで", 17, bold=True, color=DARK_NAVY)

    _add_flow_boxes(slide, Inches(3.7), [
        ("会員登録", "メール認証\n+ プロフィール"),
        ("動画視聴", "企業動画を\n体験"),
        ("エントリー", "求人応募\n/ 説明会予約"),
        ("通知受信", "メール・LINE\nリマインダー"),
    ])

    _add_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.4),
                 "審査ワークフロー（全コンテンツ共通）", 17, bold=True, color=DARK_NAVY)

    _add_flow_boxes(slide, Inches(5.6), [
        ("Draft", "下書き\n編集可能"),
        ("Submitted", "審査申請\n編集ロック"),
        ("Approved", "承認済み\n本番反映"),
        ("Active", "公開中\nON/OFF 可"),
    ], box_w=Inches(2.4), box_h=Inches(0.85))


def slide_06_video_system(prs):
    slide = _content_slide(prs, "動画配信システム", 6)

    _add_flow_boxes(slide, Inches(1.3), [
        ("S3 Upload", "元動画を保存"),
        ("MediaConvert", "HLS 変換\n1080p + 720p"),
        ("S3 HLS", "セグメント保存"),
        ("CloudFront", "CDN 配信"),
        ("HLS.js", "ブラウザ再生"),
    ], box_w=Inches(2.0), box_h=Inches(0.9))

    # Spec cards
    specs = [
        ("解像度", "1080p (5Mbps)\n720p (3Mbps)"),
        ("形式", "HLS (.m3u8 + .ts)\nH.264 + AAC"),
        ("対応入力", "MP4 / WebM\nMOV / AVI"),
        ("制限", "動画 500MB（S3 直接）\nサムネイル 5MB"),
    ]
    for i, (t, b) in enumerate(specs):
        _add_card(slide, Inches(0.8 + i * 3.1), Inches(2.7), Inches(2.9), Inches(1.2), t, b)

    # Categories
    rows = [
        ["カテゴリ", "上限", "用途", "表示場所"],
        ["メインビデオ", "1本", "企業の主力紹介動画", "企業ページトップ"],
        ["ショート動画", "10本", "社員インタビュー等", "企業ページ + トップ"],
        ["ドキュメンタリー", "10本", "企業密着・職場見学", "企業ページ + トップ"],
    ]
    _add_table(slide, Inches(0.8), Inches(4.3), Inches(11.7), rows,
               [Inches(2.2), Inches(1.0), Inches(4.0), Inches(4.5)])


def slide_07_company_page(prs):
    slide = _content_slide(prs, "企業ページ管理（Draft / Production）", 7)

    _add_flow_boxes(slide, Inches(1.3), [
        ("Draft", "企業担当者が\nスタジオで編集"),
        ("Submitted", "編集ロック\n管理者待ち"),
        ("Approved", "本番に反映\n公開 ON/OFF 可"),
        ("Active", "公開サイトに\n表示"),
    ])

    rows = [
        ["機能", "ドラフト", "本番", "審査"],
        ["企業ページ", "company_pages_draft", "company_pages", "/admin/review"],
        ["企業情報", "companies_draft", "companies", "/admin/review"],
        ["求人", "job_postings_draft", "job_postings", "/admin/jobs"],
        ["説明会", "sessions_draft", "sessions", "/admin/sessions"],
        ["動画", "videos_draft", "videos", "/admin（動画タブ）"],
    ]
    _add_table(slide, Inches(0.8), Inches(2.8), Inches(11.7), rows,
               [Inches(1.8), Inches(3.0), Inches(2.5), Inches(4.4)])

    _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
                 "admin は審査バイパスで直接入稿も可能（ドラフト approved + 本番 active を同時書き込み）",
                 13, color=MEDIUM_GRAY)


def slide_08_admin_company(prs):
    slide = _content_slide(prs, "管理者の企業管理（6タブ構成）", 8)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                 "/admin/company-accounts/[companyId] — 企業のすべてを一画面で管理",
                 15, bold=True, color=DARK_NAVY)

    tabs = [
        ("1", "企業情報", "企業名・ロゴ・業界・従業員数\nサムネイル直接更新可能"),
        ("2", "企業ページ", "キャッチコピー・紹介文・SNS\n福利厚生・カバー画像"),
        ("3", "求人", "求人一覧・新規作成\n公開 ON/OFF トグル"),
        ("4", "説明会", "説明会一覧・日程管理\n予約数の確認"),
        ("5", "動画", "動画一覧・新規追加\nカテゴリ別管理"),
        ("6", "リクルーター", "担当者アカウント管理\n招待メール送信"),
    ]

    for i, (num, title, desc) in enumerate(tabs):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(1.9 + row * 2.3)
        _add_icon_card(slide, x, y, Inches(3.8), Inches(1.9), num, title, desc)

    _add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.3),
                 "各タブの保存は admin-company-detail-actions.ts で処理。審査フローをバイパスして直接本番に反映。",
                 12, color=MEDIUM_GRAY)


def slide_09_cms(prs):
    slide = _content_slide(prs, "トップページ CMS", 9)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.35),
                 "admin が管理画面から直接 CRUD — 審査フロー不要",
                 15, bold=True, color=DARK_NAVY)

    rows = [
        ["セクション", "テーブル", "管理画面"],
        ["ヒーロースライド", "top_page_hero_items", "/admin/hero-items"],
        ["バナー", "top_page_banners", "/admin/banners"],
        ["アンバサダー", "top_page_ambassadors", "/admin/featured-videos"],
        ["就活ドキュメンタリー", "top_page_documentaries", "/admin/featured-videos"],
        ["しゅんダイアリー", "top_page_shun_diaries", "/admin/featured-videos"],
        ["就活Shorts / Videos", "top_page_featured_videos", "/admin/featured-videos"],
        ["LP サンプル動画", "lp_sample_videos", "/admin/lp-content"],
        ["LP FAQ", "lp_faq_items", "/admin/lp-content"],
        ["LP 企業ロゴ", "lp_company_logos", "/admin/lp-content"],
        ["LP スクロールバナー", "lp_scroll_banner", "/admin/lp-content"],
    ]
    _add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), rows,
               [Inches(3.0), Inches(4.5), Inches(4.2)])


def slide_10_line_and_notifications(prs):
    slide = _content_slide(prs, "LINE 連携・通知システム", 10)

    # LINE section (left)
    _add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8),
              "LINE 連携 & 配信",
              "アカウント連携:\n"
              "  LINE Login OAuth → candidates.line_user_id に保存\n"
              "  マイページで連携・解除\n\n"
              "セグメント配信 (/admin/line/broadcast):\n"
              "  卒年度・業界・職種でフィルタ → Push 送信\n"
              "  テキスト / Flex / カルーセル / 画像 / イメージマップ\n"
              "  スケジュール配信・テスト送信・自動リトライ（3回）\n"
              "  テンプレート保存・履歴閲覧・パーソナライズ変数",
              body_size=12)

    # Notification channels (right)
    _add_card(slide, Inches(6.6), Inches(1.3), Inches(5.9), Inches(2.8),
              "5 チャネル通知",
              "Email:  SendGrid REST API（14 テンプレート）\n"
              "LINE:   Messaging API Push（連携済み学生向け）\n"
              "Slack:  Incoming Webhook（運営アラート）\n"
              "Sheets: Sheets API v4（登録・予約データ記録）\n"
              "App:    DB 駆動（Studio/管理画面内通知）\n\n"
              "共通: fire-and-forget / 環境変数未設定時スキップ",
              body_size=12)

    # Notification matrix (compact)
    rows = [
        ["イベント", "Email", "LINE", "Slack", "Sheets"],
        ["会員登録", "✓", "—", "✓", "✓"],
        ["求人エントリー", "✓→学生+企業", "—", "—", "—"],
        ["説明会予約", "✓→学生+企業", "—", "—", "—"],
        ["イベント予約", "✓→学生", "✓", "✓", "✓"],
        ["イベントリマインダー", "✓→学生", "✓", "—", "—"],
        ["管理者 LINE 配信", "—", "✓（予約可）", "—", "—"],
    ]
    _add_table(slide, Inches(0.8), Inches(4.4), Inches(11.7), rows,
               [Inches(3.0), Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5)])


def slide_11_tech_stack(prs):
    slide = _content_slide(prs, "技術スタック", 11)

    # Cards layout instead of big table
    stack = [
        ("Frontend", ACCENT_BLUE,
         "Next.js 16 (App Router)\nTypeScript\nTailwind CSS 4\nshadcn/ui"),
        ("Backend / DB", DARK_NAVY,
         "Supabase (PostgreSQL)\nAuth + RLS + Realtime\nServer Actions"),
        ("動画基盤", RGBColor(0x6C, 0x5C, 0xE7),
         "Amazon S3（保存）\nMediaConvert（HLS 変換）\nCloudFront（CDN 配信）\nHLS.js（ブラウザ再生）"),
        ("インフラ", SUCCESS_GREEN,
         "Vercel（ホスティング）\nTurborepo（Monorepo）\nSendGrid（メール）\nLINE Messaging API"),
    ]
    for i, (title, bg, items) in enumerate(stack):
        x = Inches(0.6 + i * 3.2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(3.0), Inches(3.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_font(run, 18, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(12)
        run2 = p2.add_run()
        run2.text = items
        _set_font(run2, 14, color=RGBColor(0xDD, 0xEE, 0xFF))

    _add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.4),
                 "Monorepo 構成（3 アプリ）", 17, bold=True, color=DARK_NAVY)

    apps = [
        ["アプリ", "ポート", "概要"],
        ["jobtv", "3000", "動画就活情報サイト（メイン）"],
        ["event-system", "3001", "イベント運営システム"],
        ["agent-manager", "3002", "エージェント管理システム"],
    ]
    _add_table(slide, Inches(0.8), Inches(5.2), Inches(11.7), apps,
               [Inches(2.5), Inches(1.5), Inches(7.7)])


def slide_12_video_infra(prs):
    slide = _content_slide(prs, "動画インフラ詳細", 12)

    _add_flow_boxes(slide, Inches(1.3), [
        ("Upload", "元動画を S3 に保存"),
        ("MediaConvert", "2 解像度で\nHLS 変換"),
        ("S3 (HLS)", "セグメント\n+ マニフェスト"),
        ("CloudFront", "CDN\nエッジ配信"),
    ], box_w=Inches(2.6), box_h=Inches(0.9))

    # Left: bucket structure
    _add_card(slide, Inches(0.8), Inches(2.7), Inches(6.0), Inches(3.0),
              "S3 バケット構造 (jobtv-videos-stg)",
              "companies/{companyId}/videos/{draftId}/original.ext\n\n"
              "transcoded/\n"
              "  └ {landscape|portrait}/{companyId}/videos/{videoId}/hls/\n"
              "      ├ original.m3u8（マスター）\n"
              "      ├ master_1080p.m3u8 / master_720p.m3u8\n"
              "      └ *.ts（セグメント）\n\n"
              "ライフサイクル: 元動画 → 90日後 Glacier アーカイブ\n"
              "削除キューで安全な非同期削除",
              body_size=11)

    # Right: benefits + cost
    _add_card(slide, Inches(7.1), Inches(2.7), Inches(5.4), Inches(1.4),
              "アダプティブストリーミングの利点",
              "ネットワーク状況に応じた自動画質調整\nスマホ/PC/タブレット最適化\niOS Safari 標準サポート（HLS）")

    _add_card(slide, Inches(7.1), Inches(4.3), Inches(5.4), Inches(1.4),
              "コスト目安",
              "S3: 元動画の約 2 倍（2 解像度分）\nMediaConvert: 約 $0.0075/分\nCloudFront: 転送量課金")


def slide_13_db_design(prs):
    slide = _content_slide(prs, "データベース設計", 13)

    # 3 cards
    _add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(2.6),
              "Draft / Production パターン",
              "スタジオで編集 → 審査 → 公開\n\n"
              "ドラフト（*_draft）:\n"
              "  draft → submitted → approved\n\n"
              "本番テーブル:\n"
              "  承認後に反映\n"
              "  active / closed で公開制御")

    _add_card(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(2.6),
              "RLS（Row Level Security）",
              "全テーブルで RLS 有効化\n\n"
              "admin → 全件アクセス\n"
              "recruiter → 自社のみ\n"
              "candidate → 自分のみ\n"
              "anon → 公開テーブル SELECT のみ\n\n"
              "データ分離を DB レイヤーで担保")

    _add_card(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(2.6),
              "監査ログ（audit_logs）",
              "管理者・システムの操作履歴を記録\n\n"
              "action / category /\n"
              "resource_type / resource_id /\n"
              "app / metadata / ip_address\n\n"
              "カテゴリ: content_review /\n"
              "account / content_edit / access 等")

    _add_bullet_list(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.5), [
        "マイグレーション一元管理: supabase/migrations/ で全アプリ共通",
        "型自動生成: pnpm types → database.types.ts（手動編集禁止）",
        "本番 DB 変更は STG 検証済みマイグレーション経由のみ",
    ], size=13)


def slide_14_security_overview(prs):
    slide = _content_slide(prs, "セキュリティ — 多層防御", 14)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.35),
                 "11 のセキュリティレイヤーでサービスを保護",
                 15, bold=True, color=DARK_NAVY)

    # Layer 1: Network
    _add_textbox(slide, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.35),
                 "ネットワーク層", 14, bold=True, color=ACCENT_BLUE)
    _add_shield_row(slide, Inches(0.8), Inches(2.1), Inches(11.7), [
        ("DDoS 緩和", "Vercel 標準提供\nL3/L4/L7 自動遮断"),
        ("Vercel Firewall", "/admin/* を IP 制限\n許可 IP 以外は 403"),
        ("Basic 認証", "リリース前: 全体\nリリース後: /admin のみ"),
    ])

    # Layer 2: Authentication
    _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(2.5), Inches(0.35),
                 "認証・認可層", 14, bold=True, color=ACCENT_BLUE)
    _add_shield_row(slide, Inches(0.8), Inches(3.85), Inches(11.7), [
        ("Supabase Auth", "メール/パスワード認証\n全ロール共通基盤"),
        ("MFA (TOTP)", "admin 必須\nAuthenticator アプリ"),
        ("RLS", "全テーブルで有効\nロール別データ分離"),
        ("ルート保護", "requireAdmin()\nrequireStudioAuth()"),
    ])

    # Layer 3: Application
    _add_textbox(slide, Inches(0.8), Inches(5.25), Inches(2.5), Inches(0.35),
                 "アプリケーション層", 14, bold=True, color=ACCENT_BLUE)
    _add_shield_row(slide, Inches(0.8), Inches(5.6), Inches(11.7), [
        ("CAPTCHA", "Cloudflare Turnstile\nボット対策"),
        ("CSRF 防止", "state パラメータ\nCookie 照合"),
        ("Webhook 認証", "SUPABASE_HOOK_SECRET\nCRON_SECRET"),
        ("監査ログ", "操作履歴を記録\nIP・メタデータ保持"),
    ])


def slide_15_security_access(prs):
    slide = _content_slide(prs, "セキュリティ — アクセス制御詳細", 15)

    # RLS detail table
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.35),
                 "RLS ポリシーによるデータ分離", 15, bold=True, color=DARK_NAVY)

    rows = [
        ["ロール", "SELECT", "INSERT/UPDATE/DELETE"],
        ["admin", "全テーブル全件", "全テーブル全件"],
        ["recruiter", "自社 company_id のみ", "自社データのみ"],
        ["candidate", "自分の candidate_id のみ", "自分のデータのみ"],
        ["anon（未認証）", "公開テーブルのみ", "全テーブル REVOKE"],
    ]
    _add_table(slide, Inches(0.8), Inches(1.7), Inches(5.5), rows,
               [Inches(1.5), Inches(2.0), Inches(2.0)])

    # Firewall rules
    _add_textbox(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(0.35),
                 "Vercel Firewall ルール", 15, bold=True, color=DARK_NAVY)

    fw_rows = [
        ["対象パス", "種類", "条件"],
        ["/admin/*", "IP 制限", "許可 IP 以外 → 403"],
        ["/api/admin/*", "IP 制限", "API 直叩き防止"],
        ["/api/auth/login", "Rate Limit", "10 req/min per IP"],
        ["/api/studio/login", "Rate Limit", "10 req/min per IP"],
        ["/api/cron/*", "Rate Limit", "1 req/min per IP"],
    ]
    _add_table(slide, Inches(6.8), Inches(1.7), Inches(5.7), fw_rows,
               [Inches(2.0), Inches(1.5), Inches(2.2)])

    # Storage policy
    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.35),
                 "Storage ポリシー（company-assets）", 15, bold=True, color=DARK_NAVY)

    st_rows = [
        ["操作", "条件"],
        ["SELECT（公開読み取り）", "全員（public バケット）"],
        ["INSERT / UPDATE / DELETE", "admin: 全パス可\nrecruiter: 自社パスのみ"],
    ]
    _add_table(slide, Inches(0.8), Inches(4.8), Inches(5.5), st_rows,
               [Inches(2.5), Inches(3.0)])

    # DB change guard
    _add_card(slide, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.2),
              "本番 DB 変更の安全ガード",
              "マイグレーション経由のみ（直接 SQL 禁止）\n"
              "STG 検証済み → PROD 適用前バックアップ\n"
              "PITR（Point-in-Time Recovery）有効\n"
              "Claude Code hooks で危険操作を物理ブロック\n"
              "破壊的変更は 2 段階マイグレーション")


def slide_16_deploy(prs):
    slide = _content_slide(prs, "環境・デプロイ戦略", 16)

    rows = [
        ["環境", "ブランチ", "デプロイ先", "Supabase", "AWS"],
        ["ローカル", "—", "localhost:3000", "STG", "STG"],
        ["開発", "develop", "Vercel Preview", "STG", "STG"],
        ["STG", "staging", "Vercel カスタム環境", "STG", "STG"],
        ["PROD", "main", "Vercel Production", "PROD", "PROD"],
    ]
    _add_table(slide, Inches(0.8), Inches(1.3), Inches(11.7), rows,
               [Inches(1.3), Inches(1.5), Inches(3.0), Inches(3.0), Inches(2.9)])

    _add_flow_boxes(slide, Inches(3.7), [
        ("feature/*", "機能開発"),
        ("develop", "統合\n（デフォルト）"),
        ("staging", "STG\nデプロイ"),
        ("main", "PROD\nデプロイ"),
    ], box_w=Inches(2.4), box_h=Inches(0.85))

    _add_card(slide, Inches(0.8), Inches(5.0), Inches(3.7), Inches(1.3),
              "Branch Protection",
              "main: PR 必須 + 承認 1 名\nstaging: PR 必須\nforce push 禁止")

    _add_card(slide, Inches(4.8), Inches(5.0), Inches(3.7), Inches(1.3),
              "定期処理（Cron 4本）",
              "リマインダー: 毎日 JST 10:00\nLINE配信: 5分毎 / リトライ: 15分毎\nストレージ清掃: 毎日 JST 12:00")

    _add_card(slide, Inches(8.8), Inches(5.0), Inches(3.7), Inches(1.3),
              "監視・アラート",
              "Slack: メール失敗 / 新規登録 / イベント予約\nメールログ: email_logs + 管理画面")


def slide_17_roadmap(prs):
    slide = _content_slide(prs, "改善ロードマップ", 17)

    # Completed (former high priority)
    _add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.2),
              "完了済み ✓", "",
              bg=RGBColor(0xE8, 0xF8, 0xF5), title_color=SUCCESS_GREEN)
    _add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(5.3), Inches(1.5), [
        "LINE 配信履歴・ログの DB 記録 → 実装済",
        "学生向け確認メール（エントリー/予約） → 実装済",
        "LINE テスト送信・スケジュール配信・自動リトライ → 実装済",
        "LINE テンプレート保存・再利用 → 実装済",
        "S3 ストレージ管理・削除キュー → 実装済",
    ], size=13)

    # Current priorities
    _add_card(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.2),
              "高優先度", "",
              bg=RGBColor(0xFD, 0xF2, 0xE9), title_color=WARM_ORANGE)
    _add_bullet_list(slide, Inches(7.0), Inches(1.8), Inches(5.3), Inches(1.5), [
        "通知オプトアウト機能",
        "求人エントリー / 説明会の Slack・Sheets 連携",
        "Email / Slack リトライ機構",
    ], size=13)

    # Infra
    _add_card(slide, Inches(0.8), Inches(3.8), Inches(5.8), Inches(2.2),
              "インフラ・信頼性", "",
              bg=LIGHT_BLUE_BG, title_color=ACCENT_BLUE)
    _add_bullet_list(slide, Inches(1.0), Inches(4.3), Inches(5.3), Inches(1.5), [
        "Sheets 書き込み失敗の Slack 通知",
        "MediaConvert 完了の SNS 通知",
        "分散トレーシング・メトリクス拡充",
    ], size=13)

    # Future
    _add_card(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.2),
              "将来検討", "",
              bg=LIGHT_GRAY, title_color=MEDIUM_GRAY)
    _add_bullet_list(slide, Inches(7.0), Inches(4.3), Inches(5.3), Inches(1.5), [
        "LINE A/B テスト配信",
        "LINE Webhook（自動応答）",
        "全チャネル共通リトライ・キュー統合",
    ], size=13)


def slide_18_role_candidate(prs):
    slide = _content_slide(prs, "ロール別機能: 学生（Candidate）", 18)

    # Role badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "candidate"
    _set_font(run, 16, bold=True, color=WHITE, font_name=FONT_EN)

    _add_textbox(slide, Inches(3.6), Inches(1.25), Inches(8), Inches(0.4),
                 "公開ページ + /mypage  |  データ: 自分の candidate_id に紐づく情報のみ",
                 13, color=MEDIUM_GRAY)

    # Feature cards - 2 rows x 3
    features = [
        ("動画視聴", "HLS アダプティブストリーミング\n企業の動画を無料視聴\nメイン / ショート / ドキュメンタリー"),
        ("企業ページ閲覧", "企業プロフィール・求人一覧\n説明会情報の確認\n公開中（active）のみ表示"),
        ("求人エントリー", "企業ページから求人に応募\nモーダル内で完結\n重複エントリー防止（UNIQUE）"),
        ("説明会予約", "日程を選んで予約\n満席表示（RPC で件数取得）\n予約後は企業に通知メール"),
        ("LINE 連携", "マイページからワンタップ連携\nLINE Login OAuth\n連携後プッシュ通知を受信"),
        ("マイページ", "プロフィール編集\nLINE 連携状態の確認・解除\n学校マスタによるサジェスト入力"),
    ]
    for i, (title, body) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(2.0 + row * 2.3)
        _add_icon_card(slide, x, y, Inches(3.8), Inches(2.0),
                       str(i + 1), title, body)


def slide_19_role_recruiter(prs):
    slide = _content_slide(prs, "ロール別機能: 企業担当者（Recruiter）", 19)

    # Role badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "recruiter"
    _set_font(run, 16, bold=True, color=WHITE, font_name=FONT_EN)

    _add_textbox(slide, Inches(3.6), Inches(1.25), Inches(8), Inches(0.4),
                 "/studio 配下  |  データ: 自社 company_id のみ操作可能",
                 13, color=MEDIUM_GRAY)

    features = [
        ("企業ページ管理", "/studio/company\nキャッチコピー・紹介文・SNS\nカバー画像の編集・審査申請"),
        ("動画管理", "/studio/videos\nメイン / ショート / ドキュメンタリー\nアップロード・編集・審査申請"),
        ("求人管理", "/studio/jobs\n求人の作成・編集・審査申請\n公開 ON/OFF トグル"),
        ("説明会管理", "/studio/sessions\n説明会の作成・日程管理\n審査申請・公開 ON/OFF"),
        ("FAQ 管理", "/studio/faq\n企業 FAQ の作成・編集\n公開ページに表示"),
        ("設定・通知", "/studio/settings\n企業プロフィール編集\nチームメンバー招待・通知確認"),
    ]
    for i, (title, body) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(2.0 + row * 2.3)
        _add_icon_card(slide, x, y, Inches(3.8), Inches(2.0),
                       str(i + 1), title, body, icon_color=DARK_NAVY)


def slide_20_role_admin(prs):
    slide = _content_slide(prs, "ロール別機能: 管理者（Admin）", 20)

    # Role badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x6C, 0x5C, 0xE7)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "admin"
    _set_font(run, 16, bold=True, color=WHITE, font_name=FONT_EN)

    _add_textbox(slide, Inches(3.6), Inches(1.25), Inches(8), Inches(0.4),
                 "/admin 配下（MFA + IP 制限）  |  データ: 全テーブル全件アクセス可",
                 13, color=MEDIUM_GRAY)

    # Use table for admin since there are many features
    rows = [
        ["カテゴリ", "機能", "パス"],
        ["審査", "企業ページ / 企業情報 / 求人 / 説明会 / 動画の承認・却下", "/admin/review, /admin/jobs 等"],
        ["企業管理", "企業アカウント作成・6 タブ管理・審査バイパス直接入稿", "/admin/company-accounts"],
        ["学生管理", "学生アカウントの個別作成・CSV 一括登録", "/admin/student-accounts"],
        ["CMS", "ヒーロー / バナー / アンバサダー / 就活動画 / LP コンテンツ", "/admin/hero-items 等"],
        ["トップ掲載", "就活 Shorts / Videos に表示する動画の選択・並び替え", "/admin/featured-videos"],
        ["LINE 配信", "セグメント配信・スケジュール配信・テスト送信・自動リトライ", "/admin/line/broadcast"],
        ["LINE 管理", "配信履歴・テンプレート管理・リッチメニュー", "/admin/line/history 等"],
        ["メール管理", "テンプレート CRUD（14種） + 送信ログの閲覧", "/admin/email/*"],
        ["ストレージ", "S3 / Supabase Storage の削除キュー管理・孤立ファイル清掃", "/admin/storage-cleanup"],
        ["監査", "操作履歴（action / IP / metadata）の閲覧", "/admin/audit-logs"],
        ["通知", "全体通知 / 企業別通知の作成・管理", "/admin/notifications"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.9), Inches(12.7), rows,
               [Inches(1.5), Inches(6.5), Inches(4.7)])


def slide_21_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, DARK_NAVY)

    _add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
                 "まとめ", 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_accent_bar(slide, Inches(1.7), Inches(0.06))

    summaries = [
        (ACCENT_BLUE, "動画ファーストの採用体験",
         "HLS アダプティブストリーミングで\nあらゆる環境で企業のリアルを伝える\nSEO / AIO 対応で AI 時代の発見性も確保"),
        (SUCCESS_GREEN, "品質と自走の両立",
         "Draft/Production + 管理者審査で\n品質を担保しつつ企業の自走を実現\nLINE スケジュール配信・自動リトライ"),
        (RGBColor(0x6C, 0x5C, 0xE7), "堅牢なセキュリティ",
         "DDoS / RLS / MFA / IP 制限 /\nCAPTCHA / 監査ログの多層防御\nS3 削除キューで安全なストレージ管理"),
    ]

    for i, (bg, title, body) in enumerate(summaries):
        x = Inches(0.8 + i * 4.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(3.8), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_font(run, 18, bold=True, color=WHITE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        run2 = p2.add_run()
        run2.text = body
        _set_font(run2, 14, color=RGBColor(0xDD, 0xEE, 0xFF))

    _add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                 "Q & A", 28, bold=True, color=RGBColor(0xAA, 0xBB, 0xCC),
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
                 "JOBTV  |  https://media.jobtv.jp", 14, color=MEDIUM_GRAY,
                 align=PP_ALIGN.CENTER, font=FONT_EN)

    _add_footer(slide, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Section 1: Introduction (3)
    slide_01_cover(prs)               # 1
    slide_02_exec_summary(prs)        # 2
    slide_03_toc(prs)                 # 3

    # Section 2: Service & Content (4)
    slide_04_overview(prs)            # 4
    slide_05_business_flow(prs)       # 5
    slide_06_video_system(prs)        # 6
    slide_07_company_page(prs)        # 7

    # Section 3: Company Features (2)
    slide_08_admin_company(prs)       # 8
    slide_09_cms(prs)                 # 9

    # Section 4: LINE & Notifications (1)
    slide_10_line_and_notifications(prs)  # 10

    # Section 5: Tech (3)
    slide_11_tech_stack(prs)          # 11
    slide_12_video_infra(prs)         # 12
    slide_13_db_design(prs)           # 13

    # Section 6: Security (2)
    slide_14_security_overview(prs)   # 14
    slide_15_security_access(prs)     # 15

    # Section 7: Operations & Roadmap (2)
    slide_16_deploy(prs)              # 16
    slide_17_roadmap(prs)             # 17

    # Section 8: Role Features (3)
    slide_18_role_candidate(prs)      # 18
    slide_19_role_recruiter(prs)      # 19
    slide_20_role_admin(prs)          # 20

    # Summary (1)
    slide_21_summary(prs)             # 21

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"Generated: {OUTPUT_FILE}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
